# =====================================================================
# JOSH SECURITY: AUTOMATED SLIDE COMPILER (v6.0 - PRODUCTION)
# AUDITORÍA DE CONTENEDORES VISUALES Y PRESENTACIÓN TÉCNICA
# =====================================================================
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


def crear_presentacion_josh_security():
    print("[+] Configurando Motor de Diapositivas 'JOSH Security'...")
    prs = Presentation()

    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    COLOR_FONDO = RGBColor(2, 6, 23)
    COLOR_NEON = RGBColor(56, 189, 248)
    COLOR_TEXTO = RGBColor(148, 163, 184)
    COLOR_EXITO = RGBColor(16, 185, 129)

    diapositivas_data = [
        {
            "titulo": "JOSH SECURITY ENGINE: ARCHITECTURE",
            "sub": "Conexión Estratégica del Motor Python 3.13 / Flask Cloud",
            "puntos": [
                "Ecosistema Reactivo: Integración entre Flutter UI y backend.",
                "Arquitectura Asíncrona: Procesamiento en segundo plano.",
                "Motor Forense: Análisis determinista en tiempo real.",
            ],
        },
        {
            "titulo": "MECANISMO DE INGESTIÓN Y SANITIZACIÓN",
            "sub": "Protocolo de Entrada y Evaluación de Vectores de Amenaza",
            "puntos": [
                "Ingestión Multicanal: Recepción de URLs, ejecutables y "
                "teléfonos.",
                "Sanitización en Backend: Normalización bajo estándar "
                "ITU-T E.164.",
                "Fase de Inspección: Animación sincronizada del escudo 3D.",
            ],
        },
        {
            "titulo": "EVALUACIÓN DE AMENAZAS Y FIRMAS",
            "sub": "Integración Heurística con Motores Globales",
            "puntos": [
                "Análisis Multinivel: Escaneo heurístico de firmas y "
                "patrones.",
                "Consultas en Tiempo Real: Conexión con Google Safe "
                "Browsing y VT.",
                "Escala Unificada: Puntaje de 0 a 100 y niveles semánticos.",
            ],
        },
        {
            "titulo": "RESPUESTA REACTIVA Y PERSISTENCIA",
            "sub": "Mapeo de Canales y Almacenamiento Forense",
            "puntos": [
                "Mutación Dinámica de UI: Transición cromática según "
                "el riesgo.",
                "Retroalimentación Visual: Indicadores Glowing reactivos.",
                "Persistencia: Registro atómico SQLite y exportación en PDF.",
            ],
        },
        {
            "titulo": "INFRAESTRUCTURA CLOUD Y ESCALABILIDAD",
            "sub": "Hoja de Ruta del Despliegue Técnico",
            "puntos": [
                "Despliegue Cloud: Microservicios alojados en Render.",
                "Interceptación Móvil: Detección preventiva de llamadas "
                "spam.",
                "Auditoría Forense: Reportes PDF estructurados para "
                "incidentes.",
            ],
        },
    ]

    blank_layout = prs.slide_layouts[6]

    for data in diapositivas_data:
        slide = prs.slides.add_slide(blank_layout)

        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = COLOR_FONDO

        tx_box = slide.shapes.add_textbox(
            Inches(0.8), Inches(0.6), Inches(11.5), Inches(1.5)
        )
        tf = tx_box.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0

        p = tf.paragraphs[0]
        p.text = data["titulo"]
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.name = "Arial"
        p.font.color.rgb = COLOR_NEON

        p_sub = tf.add_paragraph()
        p_sub.text = data["sub"]
        p_sub.font.size = Pt(16)
        p_sub.font.italic = True
        p_sub.font.name = "Arial"
        p_sub.font.color.rgb = COLOR_EXITO
        p_sub.space_before = Pt(6)

        content_box = slide.shapes.add_textbox(
            Inches(0.8), Inches(2.4), Inches(11.5), Inches(4.3)
        )
        tf_content = content_box.text_frame
        tf_content.word_wrap = True
        tf_content.margin_left = (
            tf_content.margin_right
        ) = tf_content.margin_top = tf_content.margin_bottom = 0

        for i, punto in enumerate(data["puntos"]):
            p_pt = (
                tf_content.paragraphs[0]
                if i == 0
                else tf_content.add_paragraph()
            )
            p_pt.text = f"▪  {punto}"
            p_pt.font.size = Pt(20)
            p_pt.font.name = "Arial"
            p_pt.font.color.rgb = COLOR_TEXTO
            p_pt.space_before = Pt(20)
            p_pt.alignment = PP_ALIGN.LEFT

    output_filename = "JoshSecurity_Arquitectura_Presentacion.pptx"
    prs.save(output_filename)
    print(f"\n[+] PRESENTACIÓN COMPILADA: Guardada en '{output_filename}'.")


if __name__ == "__main__":
    crear_presentacion_josh_security()
